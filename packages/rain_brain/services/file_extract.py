"""Server-side text extraction for uploaded document attachments.

The frontend sends binary documents (PDF, DOCX, XLSX, PPTX) as base64. Without
extraction they were injected into the prompt as raw base64 — the model saw
gibberish and couldn't analyze the file. This module decodes and extracts plain
text so the model receives the actual content.

All heavy imports are lazy + guarded so a missing optional library degrades
gracefully (returns None) instead of crashing the chat request.
"""

import base64
import io
import logging

logger = logging.getLogger(__name__)

_PDF = {"pdf"}
_DOCX = {"docx"}
_XLSX = {"xlsx", "xlsm"}
_PPTX = {"pptx"}

MAX_EXTRACT_CHARS = 200_000


def _ext(name: str) -> str:
    return name.rsplit(".", 1)[-1].lower() if "." in name else ""


def is_extractable(name: str, mime: str) -> bool:
    """True if this looks like a binary document we should extract text from.

    Plain-text files (.txt/.md/.csv) return False — they are already text and
    the caller injects them verbatim.
    """
    ext = _ext(name)
    if ext in _PDF | _DOCX | _XLSX | _PPTX:
        return True
    m = (mime or "").lower()
    return any(
        k in m
        for k in ("pdf", "wordprocessing", "spreadsheet", "presentation", "excel", "msword")
    )


def _decode(content_b64: str) -> bytes | None:
    try:
        return base64.b64decode(content_b64)
    except Exception:
        return None


def _extract_pdf(data: bytes) -> str | None:
    try:
        from pypdf import PdfReader
    except Exception:
        logger.warning("pypdf not available — cannot extract PDF text")
        return None
    try:
        reader = PdfReader(io.BytesIO(data))
        parts = [t for page in reader.pages if (t := (page.extract_text() or "").strip())]
        return "\n\n".join(parts).strip() or None
    except Exception as e:
        logger.warning("PDF extraction failed: %s", e)
        return None


def _extract_docx(data: bytes) -> str | None:
    try:
        import docx  # python-docx
    except Exception:
        logger.warning("python-docx not available — cannot extract DOCX text")
        return None
    try:
        document = docx.Document(io.BytesIO(data))
        parts = [p.text for p in document.paragraphs if p.text.strip()]
        for table in document.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
        return "\n".join(parts).strip() or None
    except Exception as e:
        logger.warning("DOCX extraction failed: %s", e)
        return None


def _extract_xlsx(data: bytes) -> str | None:
    try:
        from openpyxl import load_workbook
    except Exception:
        logger.warning("openpyxl not available — cannot extract XLSX text")
        return None
    try:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts: list[str] = []
        for ws in wb.worksheets:
            parts.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = ["" if c is None else str(c) for c in row]
                if any(cells):
                    parts.append(" | ".join(cells))
        wb.close()
        return "\n".join(parts).strip() or None
    except Exception as e:
        logger.warning("XLSX extraction failed: %s", e)
        return None


def _extract_pptx(data: bytes) -> str | None:
    try:
        from pptx import Presentation
    except Exception:
        logger.info("python-pptx not available — skipping PPTX extraction")
        return None
    try:
        prs = Presentation(io.BytesIO(data))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"# Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs)
                        if line.strip():
                            parts.append(line)
        return "\n".join(parts).strip() or None
    except Exception as e:
        logger.warning("PPTX extraction failed: %s", e)
        return None


def extract_text(name: str, mime: str, content_b64: str) -> str | None:
    """Extract plain text from a base64-encoded document.

    Returns extracted text (truncated to MAX_EXTRACT_CHARS), or None if the
    type is unsupported, a required library is missing, or extraction failed.
    """
    ext = _ext(name)
    m = (mime or "").lower()
    data = _decode(content_b64)
    if data is None:
        return None

    text: str | None = None
    if ext in _PDF or "pdf" in m:
        text = _extract_pdf(data)
    elif ext in _DOCX or "wordprocessing" in m:
        text = _extract_docx(data)
    elif ext in _XLSX or "spreadsheet" in m or "excel" in m:
        text = _extract_xlsx(data)
    elif ext in _PPTX or "presentation" in m:
        text = _extract_pptx(data)

    if not text:
        return None
    if len(text) > MAX_EXTRACT_CHARS:
        text = text[:MAX_EXTRACT_CHARS] + "\n\n[...document truncated]"
    return text
